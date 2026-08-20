"""Rendu PowerPoint — restitution institutionnelle.

**C'est le livrable du moment de scène** (spec §8.7) : faire générer en direct, devant
l'assemblée, la présentation qu'elle est en train de regarder. Il doit donc s'ouvrir
sans réparation et se lire de loin — d'où le corps de texte volontairement grand et le
découpage strict d'une section par diapositive.

Comme pour Word, les caractères de contrôle sont purgés : ``python-pptx`` n'échappe
rien, et un fichier illisible se découvre au pire moment.
"""

from __future__ import annotations

import io

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Emu, Inches, Pt

from app.models.rapport import Document, Graphique, Tableau, TypeGraphique
from app.services.rendu import charte, enrichi
from app.services.rendu.ooxml import texte_xml_sur

# Auteur inscrit dans les métadonnées du fichier livré.
_AUTEUR = "OpenCacao — OpenLab Consulting"

# Dispositions du thème par défaut de python-pptx.
_DISPOSITION_TITRE = 0
_DISPOSITION_TITRE_CONTENU = 1
# Disposition « titre seul » : une figure ou un tableau occupe la place du texte.
_DISPOSITION_TITRE_SEUL = 5
# Disposition VIERGE : on dessine tout, plutôt que d'hériter du gabarit d'Office 2007.
_DISPOSITION_VIERGE = 6

# Une diapositive lue de loin ne tient pas les 800 caractères d'une section : on
# tronque proprement plutôt que de laisser le texte déborder hors du cadre.
CORPS_MAX = 600

# Corps de texte lisible depuis le fond d'une salle.
_TAILLE_CORPS = 16


def _propre(texte: str) -> str:
    """Retire ce que XML 1.0 refuse (contrôles, surrogates isolés, non-caractères)."""
    return enrichi.sans_balisage(texte_xml_sur(texte))


def _tronquer(texte: str) -> str:
    """Tronque un corps de section à la longueur lisible sur écran.

    Args:
        texte: Corps de la section.

    Returns:
        Le texte, suivi d'une ellipse s'il a été coupé.
    """
    propre = _propre(texte)
    if len(propre) <= CORPS_MAX:
        return propre
    return propre[: CORPS_MAX - 1].rstrip() + "…"


# Une figure occupe la diapositive sous son titre, marges comprises.
_FIGURE = (Inches(0.9), Inches(1.6), Inches(7.6), Inches(4.6))

_FORMES = {
    TypeGraphique.SECTEURS: XL_CHART_TYPE.PIE,
    TypeGraphique.BATONS: XL_CHART_TYPE.COLUMN_CLUSTERED,
    TypeGraphique.LIGNES: XL_CHART_TYPE.LINE_MARKERS,
}


def _rgb(hexa: str) -> RGBColor:
    """Convertit une couleur de la charte en couleur python-pptx."""
    return RGBColor.from_string(hexa)


# Hauteur du bandeau de titre des diapositives de contenu. Assez fin pour ne pas manger
# la surface utile, assez présent pour que l'œil trouve le titre sans le chercher.
_BANDEAU_PO = 0.22


def _aplat(diapositive, gauche, haut, largeur, hauteur, couleur: str) -> None:
    """Pose un rectangle de couleur pleine, sans contour."""
    forme = diapositive.shapes.add_shape(MSO_SHAPE.RECTANGLE, gauche, haut, largeur, hauteur)
    forme.fill.solid()
    forme.fill.fore_color.rgb = _rgb(couleur)
    forme.line.fill.background()
    forme.shadow.inherit = False


def _zone_texte(
    diapositive,
    gauche,
    haut,
    largeur,
    hauteur,
    texte: str,
    *,
    taille: int,
    couleur: str,
    gras: bool = False,
):
    """Ajoute une zone de texte à la charte du projet."""
    zone = diapositive.shapes.add_textbox(gauche, haut, largeur, hauteur)
    cadre = zone.text_frame
    cadre.word_wrap = True
    cadre.text = _propre(texte)
    for paragraphe in cadre.paragraphs:
        for run in paragraphe.runs:
            run.font.size = Pt(taille)
            run.font.bold = gras
            run.font.name = charte.POLICE_CORPS
            run.font.color.rgb = _rgb(couleur)
    return cadre


def _diapositive_habillee(presentation, titre: str):
    """Crée une diapositive vierge portant le bandeau et son titre.

    On part d'une disposition VIERGE plutôt que des gabarits de python-pptx : ceux-ci
    imposent la typographie et les couleurs d'Office 2007, qui sont précisément ce que
    l'audit reprochait au deck.
    """
    diapositive = presentation.slides.add_slide(presentation.slide_layouts[_DISPOSITION_VIERGE])
    largeur = presentation.slide_width
    _aplat(diapositive, 0, 0, largeur, Inches(_BANDEAU_PO), charte.ORANGE)
    _zone_texte(
        diapositive,
        Inches(0.6),
        Inches(0.45),
        largeur - Inches(1.2),
        Inches(0.9),
        titre,
        taille=26,
        couleur=charte.SOMBRE,
        gras=True,
    )
    return diapositive


def _habiller_le_theme(presentation) -> None:
    """Remplace la palette et la typographie d'Office 2007 par celles de la charte.

    python-pptx n'expose pas le thème : on réécrit la partie XML. C'est ce qui fait que
    les graphiques natifs, qui puisent leurs couleurs dans le thème, s'accordent au
    reste du document au lieu de ressortir en bleu Office 2007.
    """
    for partie in presentation.part.package.iter_parts():
        if not partie.partname.endswith("theme1.xml"):
            continue
        xml = partie.blob.decode("utf-8")
        for rang, defaut in enumerate(("4F81BD", "C0504D", "9BBB59", "8064A2", "4BACC6", "F79646")):
            xml = xml.replace(defaut, charte.couleur_serie(rang))
        xml = xml.replace('typeface="Cambria"', f'typeface="{charte.POLICE_TITRES}"')
        xml = xml.replace('typeface="Calibri"', f'typeface="{charte.POLICE_CORPS}"')
        partie._blob = xml.encode("utf-8")
        return


def _diapositive_graphique(presentation, graphique: Graphique) -> None:
    """Ajoute une diapositive portant une figure NATIVE (pas une image).

    Args:
        presentation: Présentation en construction.
        graphique: Figure à tracer.
    """
    diapositive = _diapositive_habillee(presentation, graphique.titre)
    donnees = CategoryChartData()
    donnees.categories = [_propre(c) for c in graphique.categories]
    donnees.add_series(_propre(graphique.unite or graphique.titre), graphique.valeurs)
    cadre = diapositive.shapes.add_chart(_FORMES[graphique.type], *_FIGURE, donnees)
    graphe = cadre.chart
    if graphique.type is TypeGraphique.SECTEURS:
        graphe.has_legend = True
        graphe.legend.position = XL_LEGEND_POSITION.RIGHT
        graphe.legend.include_in_layout = False
    else:
        graphe.has_legend = False


def _diapositive_tableau(presentation, tableau: Tableau) -> None:
    """Ajoute une diapositive portant un tableau.

    Le deck n'en rendait aucun : une étude défilait à l'écran sans jamais montrer un
    chiffre, alors que le Word les portait.

    Args:
        presentation: Présentation en construction.
        tableau: Tableau à rendre.
    """
    diapositive = _diapositive_habillee(presentation, tableau.titre)
    lignes, colonnes = len(tableau.lignes) + 1, len(tableau.entetes)
    forme = diapositive.shapes.add_table(lignes, colonnes, *_FIGURE)
    grille = forme.table
    for rang, entete in enumerate(tableau.entetes):
        grille.cell(0, rang).text = _propre(entete)
    for numero, ligne in enumerate(tableau.lignes, start=1):
        for rang, valeur in enumerate(ligne):
            grille.cell(numero, rang).text = _tronquer(_propre(valeur))


def rendu_pptx(document: Document) -> bytes:
    """Rend le document au format PowerPoint.

    Args:
        document: Document assemblé par le moteur.

    Returns:
        Les octets du fichier ``.pptx``.
    """
    presentation = Presentation()
    presentation.slide_width = Inches(charte.DIAPO_LARGEUR_PO)
    presentation.slide_height = Inches(charte.DIAPO_HAUTEUR_PO)
    _habiller_le_theme(presentation)
    proprietes = presentation.core_properties
    proprietes.author = _AUTEUR
    # python-pptx laisse « Steve Canny » — le nom de son auteur — dans ce champ.
    proprietes.last_modified_by = _AUTEUR
    proprietes.title = _propre(document.titre)
    proprietes.created = document.manifeste.genere_le
    proprietes.modified = document.manifeste.genere_le

    largeur, hauteur = presentation.slide_width, presentation.slide_height
    ouverture = presentation.slides.add_slide(presentation.slide_layouts[_DISPOSITION_VIERGE])
    # Bandeau plein sur le tiers haut : c'est ce qui fait qu'une couverture ressemble à
    # une couverture et non à une page blanche portant une phrase.
    _aplat(ouverture, 0, 0, largeur, Emu(int(hauteur * 0.42)), charte.ORANGE)
    _zone_texte(
        ouverture,
        Inches(0.9),
        Inches(1.1),
        largeur - Inches(1.8),
        Inches(1.8),
        document.titre,
        taille=40,
        couleur=charte.BLANC,
        gras=True,
    )
    if document.sous_titre:
        _zone_texte(
            ouverture,
            Inches(0.9),
            Emu(int(hauteur * 0.45)),
            largeur - Inches(1.8),
            Inches(0.6),
            document.sous_titre,
            taille=18,
            couleur=charte.SOMBRE,
        )
    if document.mention:
        # D5 : sur la PREMIÈRE diapositive, jamais reléguée en fin de deck.
        _zone_texte(
            ouverture,
            Inches(0.9),
            Emu(int(hauteur * 0.60)),
            largeur - Inches(1.8),
            Inches(1.2),
            document.mention,
            taille=12,
            couleur=charte.ORANGE,
            gras=True,
        )
    _zone_texte(
        ouverture,
        Inches(0.9),
        Emu(int(hauteur * 0.86)),
        largeur - Inches(1.8),
        Inches(0.4),
        "OpenCacao — OpenLab Consulting",
        taille=11,
        couleur=charte.GRIS,
    )

    for section in document.sections:
        diapositive = _diapositive_habillee(presentation, section.titre)
        _zone_texte(
            diapositive,
            Inches(0.9),
            Inches(1.7),
            presentation.slide_width - Inches(1.8),
            presentation.slide_height - Inches(2.4),
            _tronquer(section.corps),
            taille=_TAILLE_CORPS,
            couleur=charte.SOMBRE,
        )

    for graphique in document.graphiques:
        _diapositive_graphique(presentation, graphique)

    for tableau in document.tableaux:
        _diapositive_tableau(presentation, tableau)

    manifeste = document.manifeste
    finale = presentation.slides.add_slide(presentation.slide_layouts[_DISPOSITION_TITRE_CONTENU])
    finale.shapes.title.text = "Manifeste de génération"
    finale.placeholders[1].text_frame.text = "\n".join(
        (
            f"Modèle : {manifeste.modele} (version {manifeste.version_modele})",
            f"Version applicative : {manifeste.version_app}",
            f"Profil matériel : {manifeste.profil_materiel}",
            f"Généré le : {manifeste.genere_le.isoformat()}",
            f"Empreinte du demandeur : {manifeste.empreinte_demandeur}",
            f"Sources mobilisées : {len(manifeste.documents_rag)}",
        )
    )

    tampon = io.BytesIO()
    presentation.save(tampon)
    return tampon.getvalue()
