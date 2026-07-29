"""Tests de l'adaptateur PowerPoint — c'est le livrable du moment de scène."""

from __future__ import annotations

import io
from datetime import UTC, datetime

from pptx import Presentation

from app.models.constat import NiveauConfiance
from app.models.rapport import Affirmation, Document, Manifeste, Section
from app.services.rendu.diapositives import CORPS_MAX, rendu_pptx


def _document(mention: str = "", sections: int = 2, corps: str = "Un paragraphe.") -> Document:
    return Document(
        titre="Étude de filière — le cacao",
        sous_titre="Analyse documentée",
        sections=tuple(
            Section(
                titre=f"Section {index}",
                corps=corps,
                affirmations=(
                    Affirmation(
                        texte="La production avoisine 2,2 millions de tonnes.",
                        source="CNRA",
                        date="2025-10-01",
                        methode="rag",
                        confiance=NiveauConfiance.MOYENNE,
                        empreinte="e3b0c44298fc",
                    ),
                ),
            )
            for index in range(1, sections + 1)
        ),
        tableaux=(),
        manifeste=Manifeste(
            modele="opencacao-8b",
            version_modele="1.1.0",
            version_app="0.6.75",
            profil_materiel="cpu",
            genere_le=datetime(2026, 7, 29, 10, 0, tzinfo=UTC),
            empreinte_demandeur="a1b2c3d4e5f6",
        ),
        mention=mention,
    )


def _presentation(octets: bytes) -> Presentation:
    return Presentation(io.BytesIO(octets))


def _textes(octets: bytes) -> list[str]:
    return [
        forme.text_frame.text
        for diapositive in _presentation(octets).slides
        for forme in diapositive.shapes
        if forme.has_text_frame
    ]


def test_le_rendu_est_un_pptx_ouvrable():
    octets = rendu_pptx(_document())
    assert octets[:2] == b"PK"
    assert len(_presentation(octets).slides) > 0


def test_une_diapositive_de_titre_ouvre_la_presentation():
    assert "Étude de filière — le cacao" in _textes(rendu_pptx(_document()))


def test_une_diapositive_par_section_plus_titre_et_manifeste():
    """1 titre + 2 sections + 1 manifeste = 4."""
    assert len(_presentation(rendu_pptx(_document(sections=2))).slides) == 4


def test_chaque_section_donne_son_titre_et_son_corps():
    textes = _textes(rendu_pptx(_document()))
    assert "Section 1" in textes
    assert any("Un paragraphe." in texte for texte in textes)


def test_la_mention_figure_sur_la_diapositive_de_titre():
    """D5 : elle ne peut pas etre releguee en fin de deck."""
    presentation = _presentation(rendu_pptx(_document(mention="Document préparatoire.")))
    premiere = [
        forme.text_frame.text for forme in presentation.slides[0].shapes if forme.has_text_frame
    ]
    assert any("Document préparatoire." in texte for texte in premiere)


def test_la_derniere_diapositive_porte_le_manifeste():
    textes = _textes(rendu_pptx(_document()))
    assert any("opencacao-8b" in texte for texte in textes)
    assert any("a1b2c3d4e5f6" in texte for texte in textes)


def test_un_corps_trop_long_est_tronque_proprement():
    """Une diapositive lue de loin ne tient pas 800 caracteres."""
    textes = _textes(rendu_pptx(_document(corps="x" * 3000)))
    corps = next(texte for texte in textes if texte.startswith("x"))
    assert len(corps) <= CORPS_MAX
    assert corps.endswith("…")


def test_un_corps_court_n_est_pas_tronque():
    textes = _textes(rendu_pptx(_document(corps="Court.")))
    assert "Court." in textes


def test_un_caractere_de_controle_ne_corrompt_pas_la_presentation():
    """python-pptx n echappe rien : un \\x00 casserait le XML."""
    textes = _textes(rendu_pptx(_document(corps="Texte\x00 avec\x0b controle.")))
    assert "\x00" not in "".join(textes)
    assert any("Texte" in texte for texte in textes)


def test_un_document_d_une_seule_section_reste_valide():
    assert len(_presentation(rendu_pptx(_document(sections=1))).slides) == 3
